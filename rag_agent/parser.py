"""
RAG Agent with Database Routing - Multi-format Document & Vision OCR Parsing Engine.

Supports PDF, DOCX, XLSX, PPTX, PNG, JPG, WEBP, TXT, MD using:
- Vision LLM OCR (Groq LLaMA-3.2 Vision)
- Enhanced PIL Preprocessing + PyTesseract OCR
- LiteParse (Markdown Engine)
- pypdf/docx/openpyxl/pptx Dedicated Fallbacks
"""

from __future__ import annotations

import base64
import io
import os
import re
from typing import BinaryIO, Any


def chunk_markdown(raw_text: str) -> list[str]:
    """
    Splits text into chunks preserving Markdown section headers (#, ##) and
    ensuring Markdown tables (| Col1 | Col2 |) remain intact in a single chunk.
    """
    raw_text = raw_text.strip()
    if not raw_text:
        return []

    # Separate markdown text into blocks by blank lines
    blocks = re.split(r"\n\s*\n", raw_text)
    chunks: list[str] = []

    for block in blocks:
        block_strip = block.strip()
        if block_strip:
            chunks.append(block_strip)

    return [c.strip() for c in chunks if c.strip()]


def consolidate_chunks(
    chunks: list[str], max_chars: int = 1100, min_chars: int = 180
) -> list[str]:
    """
    Merges consecutive small blocks into larger, self-contained chunks so
    that table rows stay co-located with their header/context lines.

    Layout engines (e.g. LiteParse) often emit markdown tables split across
    many tiny blocks; retrieved fragments then lack the column headers the
    LLM needs to interpret numbers correctly. Merging consecutive blocks up
    to `max_chars` keeps related content together while staying inside
    embedding-model limits. Oversized blocks are kept intact.
    """
    merged: list[str] = []
    buf = ""
    for c in chunks:
        c = c.strip()
        if not c:
            continue
        if buf and len(buf) + len(c) + 2 <= max_chars:
            buf += "\n\n" + c
        else:
            if buf:
                merged.append(buf)
            if len(c) > max_chars:
                merged.append(c)
                buf = ""
            else:
                buf = c
    if buf:
        merged.append(buf)

    # Fold a trailing tiny chunk into the previous one
    if len(merged) > 1 and len(merged[-1]) < min_chars:
        merged[-2] = merged[-2] + "\n\n" + merged[-1]
        merged.pop()

    return merged


_YEAR_RE = re.compile(r"\b(?:19|20)\d{2}\b")


def serialize_tables(text: str, state: dict | None = None) -> str:
    """
    Append self-describing 'item: value (year)' lines for every markdown table
    found in the chunk, keeping the original text intact.

    Financial-statement tables usually carry their column headers (e.g.
    '**Note 2019 2018**') OUTSIDE the pipe table, and tables often split
    across chunk boundaries. Re-stating each row with explicit year labels
    removes column ambiguity for the generator and boosts keyword/BM25
    matching for value lookups. Year labels are tracked across lines so
    continuation fragments keep the right columns.
    """
    lines = text.split("\n")
    out_rows: list[str] = []
    col_years: list[str] = list(state.get("years", [])) if state else []
    table_buf: list[str] = []

    def flush() -> None:
        nonlocal table_buf
        for ln in table_buf:
            cells = [c.strip() for c in ln.strip().strip("|").split("|")]
            if all(set(c) <= set("-: ") for c in cells):  # separator row
                continue
            cells = [c for c in cells if c != ""]
            if len(cells) < 2:
                continue
            item, rest = cells[0], cells[1:]
            # If the trailing cells line up with known year columns, label them
            if col_years and len(rest) >= len(col_years):
                values, extra = rest[-len(col_years):], rest[:-len(col_years)]
                note = f" [Note {extra[0]}]" if extra and not _YEAR_RE.fullmatch(extra[0]) else ""
                pairs = "; ".join(f"{v} ({y})" for v, y in zip(values, col_years))
                out_rows.append(f"{item}: {pairs}{note}")
            else:
                out_rows.append(f"{item}: {'; '.join(rest)}")
        table_buf = []

    for ln in lines:
        stripped = ln.strip()
        if stripped.startswith("|"):
            table_buf.append(stripped)
            continue
        if table_buf:
            flush()
        years = _YEAR_RE.findall(stripped)
        if years and ("note" in stripped.lower() or "as at" in stripped.lower() or stripped.startswith("**")):
            col_years = years[:4]

    if table_buf:
        flush()

    if state is not None:
        state["years"] = col_years

    if not out_rows:
        return text
    return text + "\n\nKey figures:\n" + "\n".join(out_rows)


def serialize_tables_batch(texts: list[str]) -> list[str]:
    """Serialize a sequence of chunks, carrying column-year context forward
    so table continuation fragments keep their (year) labels."""
    state: dict = {}
    return [serialize_tables(t, state) for t in texts]


class DocumentParser:
    """Multi-format Document Parser with Multimodal Vision LLM & Enhanced OCR capabilities."""

    @staticmethod
    def _to_bytes(file_data: bytes | BinaryIO | str) -> bytes:
        if isinstance(file_data, bytes):
            return file_data
        if hasattr(file_data, "read"):
            content = file_data.read()
            if hasattr(file_data, "seek"):
                file_data.seek(0)
            return content
        with open(str(file_data), "rb") as f:
            return f.read()

    @classmethod
    def _ocr_image_via_vision_llm(
        cls, img_bytes: bytes, filename: str, groq_client: Any = None
    ) -> str:
        """Transcribe text-based images into clean structured Markdown using LLaMA 3.2 Vision LLM."""
        if not groq_client:
            api_key = os.getenv("GROQ_API_KEY")
            if api_key:
                try:
                    from openai import OpenAI
                    groq_client = OpenAI(
                        base_url="https://api.groq.com/openai/v1",
                        api_key=api_key,
                    )
                except Exception:
                    groq_client = None

        if not groq_client:
            return ""

        try:
            b64_str = base64.b64encode(img_bytes).decode("utf-8")
            fmt = "jpeg"
            if filename.lower().endswith(".png"):
                fmt = "png"
            elif filename.lower().endswith(".webp"):
                fmt = "webp"

            prompt = (
                "You are an enterprise OCR document parsing engine. Extract ALL text, titles, numbers, "
                "bullet points, headers, key-value pairs, and tables from this image verbatim. "
                "Format tables as Markdown tables. Output ONLY the extracted document text in Markdown format "
                "without any conversational intro, preamble, or commentary."
            )

            vision_models = ["llama-3.2-11b-vision-preview", "llama-3.2-90b-vision-preview"]
            for model_name in vision_models:
                try:
                    response = groq_client.chat.completions.create(
                        model=model_name,
                        messages=[
                            {
                                "role": "user",
                                "content": [
                                    {"type": "text", "text": prompt},
                                    {
                                        "type": "image_url",
                                        "image_url": {
                                            "url": f"data:image/{fmt};base64,{b64_str}"
                                        },
                                    },
                                ],
                            }
                        ],
                        temperature=0.1,
                    )
                    text = response.choices[0].message.content or ""
                    if text.strip():
                        return text.strip()
                except Exception:
                    continue
        except Exception:
            pass

        return ""

    @classmethod
    def _ocr_image_via_pytesseract(cls, img_bytes: bytes) -> str:
        """Perform PIL Image preprocessing (grayscale, contrast, sharpening) & Tesseract OCR."""
        try:
            from PIL import Image, ImageEnhance, ImageFilter
            import pytesseract

            img = Image.open(io.BytesIO(img_bytes))

            if img.mode not in ("L", "RGB"):
                img = img.convert("RGB")

            raw_text = pytesseract.image_to_string(img).strip()
            if raw_text:
                return raw_text

            gray_img = img.convert("L")
            enhancer = ImageEnhance.Contrast(gray_img)
            enhanced_img = enhancer.enhance(2.0)
            sharpened_img = enhanced_img.filter(ImageFilter.SHARPEN)

            enhanced_text = pytesseract.image_to_string(sharpened_img).strip()
            if enhanced_text:
                return enhanced_text
        except Exception:
            pass
        return ""

    @classmethod
    def parse_file(
        cls,
        file_data: bytes | BinaryIO | str,
        filename: str,
        groq_client: Any = None,
    ) -> tuple[list[str], str]:
        """
        Parses multi-format documents (PDF, DOCX, XLSX, PPTX, PNG, JPG, WEBP, TXT)
        and returns (chunks, parser_engine_label).
        """
        filename_lower = filename.lower()
        raw_bytes = cls._to_bytes(file_data)

        # Handle TXT & MD files directly
        if filename_lower.endswith((".txt", ".md")):
            content = raw_bytes.decode("utf-8", errors="ignore")
            return chunk_markdown(content), "Text Parser"

        # Handle Image files (PNG, JPG, JPEG, WEBP)
        if filename_lower.endswith((".png", ".jpg", ".jpeg", ".webp")):
            # Engine 1: Vision LLM OCR (Groq LLaMA 3.2 Vision)
            vision_text = cls._ocr_image_via_vision_llm(raw_bytes, filename, groq_client=groq_client)
            if vision_text:
                full_text = f"## Image Document ({filename}):\n{vision_text}"
                chunks = chunk_markdown(full_text)
                if chunks:
                    return chunks, "Groq Vision LLM Engine (LLaMA 3.2)"

            # Engine 2: LiteParse OCR
            try:
                from liteparse import LiteParse
                import tempfile

                with tempfile.NamedTemporaryFile(suffix=f"_{filename}", delete=False) as tmp:
                    tmp.write(raw_bytes)
                    tmp_path = tmp.name

                parser = LiteParse(output_format="markdown", ocr_enabled=True)
                parsed_res = parser.parse(tmp_path)
                extracted_text = parsed_res.text if hasattr(parsed_res, "text") else str(parsed_res)
                chunks = chunk_markdown(extracted_text)
                if chunks:
                    return chunks, "LiteParse (OCR Engine)"
            except Exception:
                pass

            # Engine 3: Enhanced PyTesseract OCR
            pytesseract_text = cls._ocr_image_via_pytesseract(raw_bytes)
            if pytesseract_text:
                full_text = f"## Image Document OCR ({filename}):\n{pytesseract_text}"
                chunks = chunk_markdown(full_text)
                if chunks:
                    return chunks, "PyTesseract (Enhanced OCR Engine)"

            # Engine 4: Pillow Fallback Image metadata
            try:
                from PIL import Image
                img = Image.open(io.BytesIO(raw_bytes))
                img_summary = f"![Image: {filename}](Format: {img.format}, Size: {img.size[0]}x{img.size[1]}, Mode: {img.mode})"
                return [img_summary], "Pillow (Image Info Engine)"
            except Exception as e:
                return [], f"Image Processing Failed ({e})"

        # Attempt LiteParse first for structured documents (PDF, DOCX, XLSX, PPTX)
        try:
            from liteparse import LiteParse
            import tempfile

            with tempfile.NamedTemporaryFile(suffix=f"_{filename}", delete=False) as tmp:
                tmp.write(raw_bytes)
                tmp_path = tmp.name

            parser = LiteParse(output_format="markdown", ocr_enabled=True)
            parsed_res = parser.parse(tmp_path)
            extracted_text = parsed_res.text if hasattr(parsed_res, "text") else str(parsed_res)
            chunks = chunk_markdown(extracted_text)
            if chunks:
                return chunks, "LiteParse (Markdown Engine)"
        except Exception:
            pass

        # Fallback 1: PDF via PyPDF
        if filename_lower.endswith(".pdf"):
            try:
                from pypdf import PdfReader
                stream = io.BytesIO(raw_bytes)
                reader = PdfReader(stream)
                pdf_text = "\n\n".join(page.extract_text() or "" for page in reader.pages)
                chunks = chunk_markdown(pdf_text)
                if chunks:
                    return chunks, "PyPDF (Fallback Engine)"
            except Exception as e:
                return [], f"PDF Parsing Failed ({e})"

        # Fallback 2: DOCX via python-docx
        if filename_lower.endswith(".docx"):
            try:
                import docx
                stream = io.BytesIO(raw_bytes)
                doc = docx.Document(stream)
                lines = []
                for p in doc.paragraphs:
                    if p.text.strip():
                        lines.append(p.text.strip())
                for table in doc.tables:
                    for row in table.rows:
                        row_cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                        lines.append("| " + " | ".join(row_cells) + " |")
                doc_text = "\n\n".join(lines)
                chunks = chunk_markdown(doc_text)
                return chunks, "python-docx (Fallback Engine)"
            except Exception as e:
                return [], f"DOCX Parsing Failed ({e})"

        # Fallback 3: XLSX via openpyxl
        if filename_lower.endswith((".xlsx", ".xls")):
            try:
                import openpyxl
                stream = io.BytesIO(raw_bytes)
                wb = openpyxl.load_workbook(stream, data_only=True)
                lines = []
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    lines.append(f"## Sheet: {sheet}")
                    for row in ws.iter_rows(values_only=True):
                        if any(cell is not None for cell in row):
                            row_vals = [str(cell).strip() if cell is not None else "" for cell in row]
                            lines.append("| " + " | ".join(row_vals) + " |")
                xlsx_text = "\n\n".join(lines)
                chunks = chunk_markdown(xlsx_text)
                return chunks, "openpyxl (Fallback Engine)"
            except Exception as e:
                return [], f"XLSX Parsing Failed ({e})"

        # Fallback 4: PPTX via python-pptx
        if filename_lower.endswith(".pptx"):
            try:
                import pptx
                stream = io.BytesIO(raw_bytes)
                prs = pptx.Presentation(stream)
                lines = []
                for idx, slide in enumerate(prs.slides, 1):
                    lines.append(f"## Slide {idx}")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            lines.append(shape.text.strip())
                pptx_text = "\n\n".join(lines)
                chunks = chunk_markdown(pptx_text)
                return chunks, "python-pptx (Fallback Engine)"
            except Exception as e:
                return [], f"PPTX Parsing Failed ({e})"

        return [], "Unknown File Format"
