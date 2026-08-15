"""
Unit tests for DocumentParser and Markdown section chunking.
"""

from rag_agent.parser import DocumentParser, chunk_markdown


def test_chunk_markdown_preserves_tables_and_headers():
    markdown_text = """
# Product Specifications

Here is the introduction to TechPro X1.

| Model | RAM | Storage | Price |
| --- | --- | --- | --- |
| TechPro X1 | 16GB | 512GB | $999 |
| TechPro Ultra | 32GB | 1TB | $1499 |

## Support Guidelines
Contact customer support for pass reset.
"""
    chunks = chunk_markdown(markdown_text)

    # Verify tables and headers are parsed as distinct intact chunks
    assert len(chunks) >= 2
    # Ensure table block remains contiguous
    table_chunk = [c for c in chunks if "|" in c][0]
    assert "| Model | RAM | Storage | Price |" in table_chunk
    assert "| TechPro Ultra | 32GB | 1TB | $1499 |" in table_chunk


def test_document_parser_txt_file():
    txt_content = b"First section.\n\nSecond section with details."
    chunks, engine = DocumentParser.parse_file(txt_content, "sample.txt")

    assert len(chunks) == 2
    assert chunks[0] == "First section."
    assert "Text Parser" in engine


def test_document_parser_docx_fallback():
    import io
    import docx
    doc = docx.Document()
    doc.add_paragraph("Paragraph inside docx document.")
    table = doc.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Header 1"
    table.cell(0, 1).text = "Header 2"

    stream = io.BytesIO()
    doc.save(stream)
    docx_bytes = stream.getvalue()

    chunks, engine = DocumentParser.parse_file(docx_bytes, "sample.docx")
    assert len(chunks) > 0
    assert "docx" in engine.lower() or "liteparse" in engine.lower()


def test_document_parser_xlsx_fallback():
    import io
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "DataSheet"
    ws.append(["Item", "Price"])
    ws.append(["Laptop", "$999"])

    stream = io.BytesIO()
    wb.save(stream)
    xlsx_bytes = stream.getvalue()

    chunks, engine = DocumentParser.parse_file(xlsx_bytes, "data.xlsx")
    assert len(chunks) > 0
    assert "openpyxl" in engine.lower() or "liteparse" in engine.lower()


def test_document_parser_pptx_fallback():
    import io
    import pptx
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Presentation Title"

    stream = io.BytesIO()
    prs.save(stream)
    pptx_bytes = stream.getvalue()

    chunks, engine = DocumentParser.parse_file(pptx_bytes, "deck.pptx")
    assert len(chunks) > 0
    assert "pptx" in engine.lower() or "liteparse" in engine.lower()


def test_document_parser_image_info():
    import io
    from PIL import Image
    img = Image.new("RGB", (100, 100), color="blue")
    stream = io.BytesIO()
    img.save(stream, format="PNG")
    img_bytes = stream.getvalue()

    chunks, engine = DocumentParser.parse_file(img_bytes, "logo.png")
    assert len(chunks) > 0
    assert any(k in engine.lower() for k in ["pillow", "liteparse", "ocr", "image", "vision"])


def test_document_parser_vision_llm_mock():
    import io
    from unittest.mock import MagicMock
    from PIL import Image

    img = Image.new("RGB", (200, 200), color="white")
    stream = io.BytesIO()
    img.save(stream, format="PNG")
    img_bytes = stream.getvalue()

    mock_groq_client = MagicMock()
    mock_choice = MagicMock()
    mock_choice.message.content = "## Extracted Text:\nProduct Specs: TechPro Laptop\nPrice: $999"
    mock_response = MagicMock()
    mock_response.choices = [mock_choice]
    mock_groq_client.chat.completions.create.return_value = mock_response

    chunks, engine = DocumentParser.parse_file(img_bytes, "specs.png", groq_client=mock_groq_client)
    assert len(chunks) > 0
    assert "TechPro Laptop" in "\n".join(chunks)
    assert "Vision LLM Engine" in engine


def test_document_parser_text_image_ocr():
    import io
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (400, 100), color=(255, 255, 255))
    d = ImageDraw.Draw(img)
    d.text((10, 10), "TechPro Laptop OCR Test", fill=(0, 0, 0))

    stream = io.BytesIO()
    img.save(stream, format="PNG")
    img_bytes = stream.getvalue()

    chunks, engine = DocumentParser.parse_file(img_bytes, "test_ocr.png")
    assert len(chunks) > 0


